from __future__ import annotations

from pathlib import Path
import tempfile
import unittest

from pure_ai_chat_loader import AI_CHAT_ROOT, load_document_artifacts_module


class DocumentArtifactTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls) -> None:
        cls.module = load_document_artifacts_module()

    def test_txt_artifact_is_utf8_fixed_extension_and_verified(self):
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            result = self.module.create_document_artifact(
                "create_txt_document",
                "开发记录",
                "第一项\n第二项",
                workspace_root=root,
            )

            target = root / result.relative_file
            self.assertEqual(target.suffix, ".txt")
            self.assertTrue(target.name.startswith("artifact_"))
            self.assertNotIn("开发记录", target.name)
            self.assertEqual(result.format_label, "TXT")
            self.assertEqual(result.item_count, len("第一项\n第二项"))
            self.assertEqual(
                target.read_text(encoding="utf-8"),
                "开发记录\n====\n\n第一项\n第二项\n",
            )
            self.assertEqual(result.sha256[:12], result.short_sha256)
            self.assertFalse(list(root.glob(".document-artifact-*")))

    def test_word_artifact_is_valid_docx_with_headings_and_lists(self):
        from docx import Document

        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            result = self.module.create_document_artifact(
                "create_word_document",
                "AIchatbot 开发报告",
                "# 当前状态\n正文内容。\n## 已完成\n- 文本聊天\n- 本地表情",
                workspace_root=root,
            )

            target = root / result.relative_file
            document = Document(target)
            texts = [paragraph.text for paragraph in document.paragraphs]
            self.assertEqual(target.suffix, ".docx")
            self.assertEqual(document.core_properties.title, "AIchatbot 开发报告")
            self.assertIn("当前状态", texts)
            self.assertIn("文本聊天", texts)
            self.assertGreaterEqual(result.item_count, 6)
            self.assertFalse(list(root.glob(".document-artifact-*")))

    def test_presentation_artifact_uses_heading_slide_boundaries(self):
        from pptx import Presentation

        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            result = self.module.create_document_artifact(
                "create_presentation",
                "MainAgent 能力规划",
                "## 当前能力\n- 只读诊断\n- 审批写入\n## 下一阶段\n- TXT\n- Word\n- PPT",
                workspace_root=root,
            )

            target = root / result.relative_file
            presentation = Presentation(target)
            titles = [slide.shapes.title.text for slide in presentation.slides]
            self.assertEqual(target.suffix, ".pptx")
            self.assertEqual(len(presentation.slides), 3)
            self.assertEqual(result.item_count, 3)
            self.assertEqual(titles, ["MainAgent 能力规划", "当前能力", "下一阶段"])
            title_run = presentation.slides[0].shapes.title.text_frame.paragraphs[0].runs[0]
            section_run = presentation.slides[1].shapes.title.text_frame.paragraphs[0].runs[0]
            self.assertGreaterEqual(title_run.font.size.pt, 50)
            self.assertGreaterEqual(section_run.font.size.pt, 35)
            slide_texts = [
                shape.text
                for shape in presentation.slides[1].shapes
                if hasattr(shape, "text_frame") and shape.has_text_frame
            ]
            self.assertTrue(any("AIchatbot · MainAgent" in text for text in slide_texts))
            self.assertTrue(any("• 只读诊断" in text for text in slide_texts))
            body_shape = next(
                shape
                for shape in presentation.slides[1].shapes
                if hasattr(shape, "text") and "• 只读诊断" in shape.text
            )
            body_sizes = [
                run.font.size.pt
                for paragraph in body_shape.text_frame.paragraphs
                for run in paragraph.runs
                if run.text.strip()
            ]
            self.assertTrue(body_sizes)
            self.assertGreaterEqual(min(body_sizes), 20)
            self.assertFalse(list(root.glob(".document-artifact-*")))

    def test_delivery_revalidates_hash_and_rejects_tampering(self):
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            result = self.module.create_document_artifact(
                "create_txt_document",
                "发送测试",
                "原始内容",
                workspace_root=root,
            )
            delivery = self.module.prepare_document_artifact_delivery(
                result,
                workspace_root=root,
            )
            verified = self.module.validate_document_artifact_delivery(
                delivery,
                workspace_root=root,
            )
            self.assertEqual(verified.sha256, result.sha256)
            (root / result.relative_file).write_text("篡改内容", encoding="utf-8")
            with self.assertRaises(self.module.DocumentArtifactError) as raised:
                self.module.validate_document_artifact_delivery(
                    delivery,
                    workspace_root=root,
                )
            self.assertEqual(raised.exception.code, "delivery_integrity_failed")

    def test_invalid_commands_text_budgets_and_slide_count_fail_closed(self):
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            cases = (
                ("unknown", "标题", "内容", "unsupported_document_command"),
                ("create_txt_document", "", "内容", "invalid_title"),
                ("create_txt_document", "标题\n换行", "内容", "invalid_title"),
                ("create_txt_document", "标题", "", "invalid_content"),
                ("create_txt_document", "标题", "包含\x00控制符", "invalid_content"),
            )
            for command, title, content, code in cases:
                with self.subTest(code=code):
                    with self.assertRaises(self.module.DocumentArtifactError) as raised:
                        self.module.create_document_artifact(
                            command,
                            title,
                            content,
                            workspace_root=root,
                        )
                    self.assertEqual(raised.exception.code, code)

            too_many_slides = "\n".join(
                f"## 幻灯片 {index}\n内容" for index in range(20)
            )
            with self.assertRaises(self.module.DocumentArtifactError) as raised:
                self.module.create_document_artifact(
                    "create_presentation",
                    "过多幻灯片",
                    too_many_slides,
                    workspace_root=root,
                )
            self.assertEqual(raised.exception.code, "too_many_slides")
            self.assertEqual(list(root.iterdir()), [])

    def test_presentation_slide_count_matches_heading_and_continuation_rules(self):
        content = "\n".join(
            [
                "# 重复的演示文稿标题",
                "## 能力概览",
                *[f"- 能力 {index}" for index in range(9)],
                "## 下一步",
                "- 计划",
            ]
        )

        # Generated title slide + two continuation slides for the first section
        # + one slide for the second section.
        self.assertEqual(self.module.presentation_slide_count(content), 4)

    def test_workspace_symbolic_link_is_rejected_when_supported(self):
        with tempfile.TemporaryDirectory() as directory:
            base = Path(directory)
            real_root = base / "real"
            real_root.mkdir()
            linked_root = base / "linked"
            try:
                linked_root.symlink_to(real_root, target_is_directory=True)
            except OSError:
                self.skipTest("symbolic links are unavailable in this environment")
            with self.assertRaises(self.module.DocumentArtifactError) as raised:
                self.module.create_document_artifact(
                    "create_txt_document",
                    "标题",
                    "内容",
                    workspace_root=linked_root,
                )
            self.assertEqual(raised.exception.code, "workspace_link_rejected")

    def test_module_has_no_network_database_nonebot_or_shell_dependency(self):
        source = (AI_CHAT_ROOT / "document_artifacts.py").read_text(encoding="utf-8").lower()
        for forbidden in (
            "nonebot",
            "httpx",
            "openai",
            "tavily",
            "sqlite",
            "database",
            "subprocess",
            "powershell",
            "send_private_msg",
        ):
            self.assertNotIn(forbidden, source)

        plugin_source = (AI_CHAT_ROOT / "__init__.py").read_text(encoding="utf-8")
        executor_start = plugin_source.index("def execute_owner_write_command(")
        executor_end = plugin_source.index(
            "def execute_document_delivery_command(",
            executor_start,
        )
        executor = plugin_source[executor_start:executor_end]
        self.assertIn("create_document_artifact(", executor)
        self.assertIn("未通过 QQ 发送文件", executor)
        self.assertNotIn("send_private_msg", executor)

        delivery_start = plugin_source.index("async def _send_new_document_deliveries(")
        delivery_end = plugin_source.index(
            "def create_main_agent_approval_resume_tool_registry(",
            delivery_start,
        )
        delivery_executor = plugin_source[delivery_start:delivery_end]
        self.assertEqual(delivery_executor.count("await bot.call_api("), 1)
        self.assertIn('"send_private_msg"', delivery_executor)
        self.assertIn('MessageSegment(\n                "file",', delivery_executor)
        self.assertIn("validate_document_artifact_delivery(", delivery_executor)
        self.assertIn("_pending_document_deliveries.pop(", delivery_executor)
        self.assertIn("PrivateMessageEvent", delivery_executor)
        self.assertIn("is_owner(config, event)", delivery_executor)
        self.assertNotIn("while ", delivery_executor)
        self.assertNotIn("for attempt", delivery_executor)

        pyproject = (AI_CHAT_ROOT.parents[2] / "pyproject.toml").read_text(encoding="utf-8")
        self.assertIn('"python-docx>=1.2,<2"', pyproject)
        self.assertIn('"python-pptx>=1.0,<2"', pyproject)
        gitignore = (AI_CHAT_ROOT.parents[2] / ".gitignore").read_text(encoding="utf-8")
        self.assertIn("output/main-agent-workspace/", gitignore)


if __name__ == "__main__":
    unittest.main()
