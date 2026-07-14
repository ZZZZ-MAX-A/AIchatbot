from __future__ import annotations

from dataclasses import dataclass
from datetime import datetime, timezone
import hashlib
import os
from pathlib import Path
import re
import secrets
import tempfile
import zipfile


DOCUMENT_ARTIFACT_COMMANDS = (
    "create_txt_document",
    "create_word_document",
    "create_presentation",
)
DOCUMENT_ARTIFACT_MAX_TITLE_CHARS = 120
DOCUMENT_ARTIFACT_MAX_CONTENT_CHARS = 20_000
DOCUMENT_ARTIFACT_MAX_BYTES = 10 * 1024 * 1024
DOCUMENT_ARTIFACT_MAX_SLIDES = 20

_COMMAND_EXTENSIONS = {
    "create_txt_document": ".txt",
    "create_word_document": ".docx",
    "create_presentation": ".pptx",
}
_COMMAND_LABELS = {
    "create_txt_document": "TXT",
    "create_word_document": "Word",
    "create_presentation": "PowerPoint",
}


class DocumentArtifactError(RuntimeError):
    def __init__(self, code: str) -> None:
        super().__init__(code)
        self.code = code


@dataclass(frozen=True)
class DocumentArtifactResult:
    artifact_id: str
    command: str
    format_label: str
    relative_file: str
    bytes: int
    sha256: str
    item_count: int

    @property
    def short_sha256(self) -> str:
        return self.sha256[:12]


@dataclass(frozen=True)
class DocumentArtifactDelivery:
    artifact_id: str
    command: str
    file_path: Path
    bytes: int
    sha256: str


def default_document_artifact_root() -> Path:
    return Path(__file__).resolve().parents[3] / "output" / "main-agent-workspace"


def _validated_title(value: object) -> str:
    if not isinstance(value, str):
        raise DocumentArtifactError("invalid_title")
    title = value.strip()
    if not title or len(title) > DOCUMENT_ARTIFACT_MAX_TITLE_CHARS:
        raise DocumentArtifactError("invalid_title")
    if "\n" in title or "\r" in title or _has_disallowed_control(title):
        raise DocumentArtifactError("invalid_title")
    return title


def _validated_content(value: object) -> str:
    if not isinstance(value, str):
        raise DocumentArtifactError("invalid_content")
    content = value.replace("\r\n", "\n").replace("\r", "\n").strip()
    if not content or len(content) > DOCUMENT_ARTIFACT_MAX_CONTENT_CHARS:
        raise DocumentArtifactError("invalid_content")
    if _has_disallowed_control(content):
        raise DocumentArtifactError("invalid_content")
    return content


def _has_disallowed_control(value: str) -> bool:
    return any(ord(character) < 32 and character not in "\n\t" for character in value)


def _prepare_workspace(root: Path, *, enforce_fixed_root: bool) -> Path:
    project_root = Path(__file__).resolve().parents[3]
    candidate = root
    if enforce_fixed_root:
        fixed_root = default_document_artifact_root()
        if candidate != fixed_root:
            raise DocumentArtifactError("workspace_outside_fixed_root")
        for path in (project_root / "output", fixed_root):
            if path.exists() and path.is_symlink():
                raise DocumentArtifactError("workspace_link_rejected")
    elif candidate.exists() and candidate.is_symlink():
        raise DocumentArtifactError("workspace_link_rejected")

    try:
        candidate.mkdir(parents=True, exist_ok=True)
    except OSError as exc:
        raise DocumentArtifactError("workspace_unavailable") from exc
    if not candidate.is_dir() or candidate.is_symlink():
        raise DocumentArtifactError("workspace_invalid")

    resolved = candidate.resolve(strict=True)
    if enforce_fixed_root:
        try:
            resolved.relative_to(project_root.resolve(strict=True))
        except ValueError as exc:
            raise DocumentArtifactError("workspace_outside_project") from exc
    return resolved


def _next_artifact_target(root: Path, extension: str) -> tuple[str, Path]:
    for _ in range(20):
        stamp = datetime.now(timezone.utc).strftime("%Y%m%dT%H%M%SZ")
        artifact_id = f"artifact_{stamp}_{secrets.token_hex(4)}"
        target = root / f"{artifact_id}{extension}"
        if not target.exists():
            return artifact_id, target
    raise DocumentArtifactError("artifact_id_unavailable")


def _temporary_path(root: Path, extension: str) -> Path:
    file_descriptor, name = tempfile.mkstemp(
        prefix=".document-artifact-",
        suffix=extension,
        dir=root,
    )
    os.close(file_descriptor)
    return Path(name)


def _fsync_file(path: Path) -> None:
    with path.open("r+b") as file_handle:
        os.fsync(file_handle.fileno())


def _validate_package(path: Path, required_member: str) -> None:
    try:
        with zipfile.ZipFile(path, "r") as archive:
            if archive.testzip() is not None:
                raise DocumentArtifactError("package_corrupt")
            names = set(archive.namelist())
    except (OSError, zipfile.BadZipFile) as exc:
        raise DocumentArtifactError("package_corrupt") from exc
    if "[Content_Types].xml" not in names or required_member not in names:
        raise DocumentArtifactError("package_invalid")


def _render_txt(path: Path, title: str, content: str) -> int:
    payload = f"{title}\n{'=' * len(title)}\n\n{content}\n"
    try:
        with path.open("w", encoding="utf-8", newline="\n") as file_handle:
            file_handle.write(payload)
            file_handle.flush()
            os.fsync(file_handle.fileno())
        if path.read_text(encoding="utf-8") != payload:
            raise DocumentArtifactError("txt_verify_failed")
    except OSError as exc:
        raise DocumentArtifactError("txt_write_failed") from exc
    return len(content)


def _set_docx_style_font(style, font_name: str) -> None:
    from docx.oxml.ns import qn

    style.font.name = font_name
    style._element.get_or_add_rPr().rFonts.set(qn("w:eastAsia"), font_name)


def _render_docx(path: Path, title: str, content: str) -> int:
    try:
        from docx import Document
        from docx.shared import Pt
    except ImportError as exc:
        raise DocumentArtifactError("docx_dependency_unavailable") from exc

    document = Document()
    document.core_properties.title = title
    document.core_properties.author = "AIchatbot MainAgent"
    normal_style = document.styles["Normal"]
    _set_docx_style_font(normal_style, "Microsoft YaHei")
    normal_style.font.size = Pt(11)
    for style_name in ("Title", "Heading 1", "Heading 2", "Heading 3"):
        _set_docx_style_font(document.styles[style_name], "Microsoft YaHei")
    document.add_heading(title, level=0)

    numbered_pattern = re.compile(r"^\d+[.)、]\s*(.+)$")
    for raw_line in content.splitlines():
        stripped = raw_line.strip()
        if not stripped:
            continue
        if stripped.startswith("### "):
            document.add_heading(stripped[4:].strip(), level=3)
        elif stripped.startswith("## "):
            document.add_heading(stripped[3:].strip(), level=2)
        elif stripped.startswith("# "):
            document.add_heading(stripped[2:].strip(), level=1)
        elif stripped.startswith(("- ", "* ", "• ")):
            document.add_paragraph(stripped[2:].strip(), style="List Bullet")
        elif match := numbered_pattern.match(stripped):
            document.add_paragraph(match.group(1).strip(), style="List Number")
        else:
            document.add_paragraph(stripped)

    try:
        document.save(path)
        _fsync_file(path)
        _validate_package(path, "word/document.xml")
        verified = Document(path)
    except DocumentArtifactError:
        raise
    except Exception as exc:
        raise DocumentArtifactError("docx_write_failed") from exc
    paragraph_count = sum(1 for paragraph in verified.paragraphs if paragraph.text.strip())
    if paragraph_count < 2:
        raise DocumentArtifactError("docx_verify_failed")
    return paragraph_count


def _ppt_sections(content: str) -> list[tuple[str, list[str]]]:
    sections: list[tuple[str, list[str]]] = []
    current_title = ""
    current_lines: list[str] = []

    def flush() -> None:
        nonlocal current_title, current_lines
        lines = [line for line in current_lines if line.strip()]
        if current_title or lines:
            sections.append((current_title.strip(), lines))
        current_title = ""
        current_lines = []

    for raw_line in content.splitlines():
        stripped = raw_line.strip()
        if stripped == "---":
            flush()
            continue
        if stripped.startswith("## "):
            flush()
            current_title = stripped[3:].strip()
            continue
        if stripped.startswith("# ") and not sections and not current_title:
            # The deck title is supplied separately and rendered on the generated
            # title slide. Ignore a repeated Markdown H1 instead of creating a
            # duplicate first content slide.
            continue
        current_lines.append(raw_line)
    flush()

    if not sections:
        lines = [line.strip() for line in content.splitlines() if line.strip()]
        sections = [("内容", lines)]

    expanded: list[tuple[str, list[str]]] = []
    for index, (section_title, lines) in enumerate(sections, start=1):
        clean_lines = [
            re.sub(
                r"^(?:#{1,6}\s+|[-*•]\s*|\d+[.)、]\s*)",
                "",
                line.strip(),
            )
            for line in lines
            if line.strip()
        ]
        chunks = [clean_lines[offset : offset + 8] for offset in range(0, len(clean_lines), 8)]
        if not chunks:
            chunks = [[]]
        for chunk_index, chunk in enumerate(chunks, start=1):
            title = section_title or f"内容 {index}"
            if len(chunks) > 1:
                title = f"{title}（{chunk_index}/{len(chunks)}）"
            expanded.append((title, chunk))
    return expanded


def presentation_slide_count(content: object) -> int:
    """Return the exact rendered slide count, including the generated title slide."""
    safe_content = _validated_content(content)
    return len(_ppt_sections(safe_content)) + 1


def _set_ppt_background(slide, rgb: tuple[int, int, int]) -> None:
    from pptx.dml.color import RGBColor

    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)


def _add_ppt_rect(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    rgb: tuple[int, int, int],
) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*rgb)
    shape.line.fill.background()


def _add_ppt_textbox(
    slide,
    text: str,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int,
    rgb: tuple[int, int, int],
    bold: bool = False,
    alignment=None,
    vertical_anchor=None,
):
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    text_frame = box.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    if vertical_anchor is not None:
        text_frame.vertical_anchor = vertical_anchor
    paragraph = text_frame.paragraphs[0]
    if alignment is not None:
        paragraph.alignment = alignment
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*rgb)
    return box


def _add_ppt_footer(slide, slide_number: int, accent: tuple[int, int, int]) -> None:
    from pptx.enum.text import PP_ALIGN

    _add_ppt_textbox(
        slide,
        "AIchatbot · MainAgent",
        left=0.8,
        top=7.08,
        width=3.5,
        height=0.2,
        font_size=10,
        rgb=(100, 116, 139),
    )
    _add_ppt_textbox(
        slide,
        f"{slide_number:02d}",
        left=11.8,
        top=7.0,
        width=0.7,
        height=0.28,
        font_size=12,
        rgb=accent,
        bold=True,
        alignment=PP_ALIGN.RIGHT,
    )


def _render_pptx(path: Path, title: str, content: str) -> int:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise DocumentArtifactError("pptx_dependency_unavailable") from exc

    sections = _ppt_sections(content)
    if presentation_slide_count(content) > DOCUMENT_ARTIFACT_MAX_SLIDES:
        raise DocumentArtifactError("too_many_slides")

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    presentation.core_properties.title = title
    presentation.core_properties.author = "AIchatbot MainAgent"

    navy = (15, 23, 42)
    slate = (51, 65, 85)
    muted = (100, 116, 139)
    white = (248, 250, 252)
    surface = (244, 247, 251)
    accents = (
        (37, 99, 235),
        (13, 148, 136),
        (124, 58, 237),
    )

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    _set_ppt_background(title_slide, navy)
    _add_ppt_rect(
        title_slide,
        left=0.9,
        top=1.45,
        width=1.35,
        height=0.08,
        rgb=accents[0],
    )
    title_shape = title_slide.shapes.title
    title_shape.left = Inches(0.9)
    title_shape.top = Inches(1.85)
    title_shape.width = Inches(11.2)
    title_shape.height = Inches(2.05)
    title_shape.text = title
    title_shape.text_frame.word_wrap = True
    title_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_shape.text_frame.margin_left = 0
    title_shape.text_frame.margin_right = 0
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(52)
            run.font.bold = True
            run.font.color.rgb = RGBColor(*white)
    subtitle_shape = title_slide.placeholders[1]
    subtitle_shape.left = Inches(0.95)
    subtitle_shape.top = Inches(4.25)
    subtitle_shape.width = Inches(8.5)
    subtitle_shape.height = Inches(0.5)
    subtitle_shape.text = "AIchatbot MainAgent · 受控文档产物"
    subtitle_shape.text_frame.margin_left = 0
    subtitle_shape.text_frame.margin_right = 0
    for paragraph in subtitle_shape.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor(148, 163, 184)
    _add_ppt_textbox(
        title_slide,
        "01",
        left=11.2,
        top=6.55,
        width=1.0,
        height=0.42,
        font_size=16,
        rgb=accents[0],
        bold=True,
        alignment=PP_ALIGN.RIGHT,
    )

    for slide_index, (section_title, lines) in enumerate(sections, start=2):
        accent = accents[(slide_index - 2) % len(accents)]
        is_closing = any(
            marker in section_title
            for marker in ("下一步", "总结", "结论", "展望", "行动")
        )
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        _set_ppt_background(slide, (238, 246, 246) if is_closing else surface)
        _add_ppt_rect(
            slide,
            left=0,
            top=0,
            width=13.333,
            height=0.1,
            rgb=accent,
        )
        _add_ppt_textbox(
            slide,
            f"{slide_index:02d}",
            left=0.85,
            top=0.55,
            width=0.65,
            height=0.5,
            font_size=18,
            rgb=accent,
            bold=True,
        )
        section_title_shape = slide.shapes.title
        section_title_shape.left = Inches(1.55)
        section_title_shape.top = Inches(0.48)
        section_title_shape.width = Inches(10.6)
        section_title_shape.height = Inches(0.72)
        section_title_shape.text = section_title
        section_title_shape.text_frame.word_wrap = True
        section_title_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        section_title_shape.text_frame.margin_left = 0
        section_title_shape.text_frame.margin_right = 0
        for paragraph in section_title_shape.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.name = "Microsoft YaHei"
                run.font.size = Pt(36)
                run.font.bold = True
                run.font.color.rgb = RGBColor(*navy)
        _add_ppt_rect(
            slide,
            left=0.9,
            top=1.35,
            width=11.55,
            height=0.015,
            rgb=(203, 213, 225),
        )

        body_box = slide.shapes.add_textbox(
            Inches(1.1),
            Inches(1.75),
            Inches(11.0),
            Inches(4.85),
        )
        text_frame = body_box.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.08)
        text_frame.margin_right = Inches(0.08)
        text_frame.margin_top = Inches(0.08)
        text_frame.margin_bottom = Inches(0.08)
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        if lines:
            for index, line in enumerate(lines):
                paragraph = (
                    text_frame.paragraphs[0]
                    if index == 0
                    else text_frame.add_paragraph()
                )
                paragraph.text = f"• {line}"
                paragraph.level = 0
                paragraph.alignment = PP_ALIGN.LEFT
                paragraph.space_after = Pt(13)
                paragraph.line_spacing = 1.12
                for run in paragraph.runs:
                    run.font.name = "Microsoft YaHei"
                    run.font.size = Pt(22)
                    run.font.color.rgb = RGBColor(*slate)
        else:
            text_frame.paragraphs[0].text = ""
        _add_ppt_footer(slide, slide_index, accent)

    try:
        presentation.save(path)
        _fsync_file(path)
        _validate_package(path, "ppt/presentation.xml")
        verified = Presentation(path)
    except DocumentArtifactError:
        raise
    except Exception as exc:
        raise DocumentArtifactError("pptx_write_failed") from exc
    if len(verified.slides) != len(sections) + 1:
        raise DocumentArtifactError("pptx_verify_failed")
    return len(verified.slides)


def create_document_artifact(
    command: str,
    title: object,
    content: object,
    *,
    workspace_root: Path | None = None,
) -> DocumentArtifactResult:
    if command not in DOCUMENT_ARTIFACT_COMMANDS:
        raise DocumentArtifactError("unsupported_document_command")
    safe_title = _validated_title(title)
    safe_content = _validated_content(content)
    root = _prepare_workspace(
        workspace_root or default_document_artifact_root(),
        enforce_fixed_root=workspace_root is None,
    )
    extension = _COMMAND_EXTENSIONS[command]
    artifact_id, target = _next_artifact_target(root, extension)
    temporary = _temporary_path(root, extension)
    try:
        if command == "create_txt_document":
            item_count = _render_txt(temporary, safe_title, safe_content)
        elif command == "create_word_document":
            item_count = _render_docx(temporary, safe_title, safe_content)
        else:
            item_count = _render_pptx(temporary, safe_title, safe_content)
        size = temporary.stat().st_size
        if size <= 0 or size > DOCUMENT_ARTIFACT_MAX_BYTES:
            raise DocumentArtifactError("artifact_size_invalid")
        digest = hashlib.sha256(temporary.read_bytes()).hexdigest()
        os.replace(temporary, target)
    except Exception:
        try:
            temporary.unlink(missing_ok=True)
        except OSError:
            pass
        raise

    if target.parent.resolve(strict=True) != root:
        target.unlink(missing_ok=True)
        raise DocumentArtifactError("artifact_path_invalid")
    final_payload = target.read_bytes()
    final_digest = hashlib.sha256(final_payload).hexdigest()
    if len(final_payload) != size or final_digest != digest:
        target.unlink(missing_ok=True)
        raise DocumentArtifactError("artifact_final_verify_failed")
    return DocumentArtifactResult(
        artifact_id=artifact_id,
        command=command,
        format_label=_COMMAND_LABELS[command],
        relative_file=(
            target.relative_to(Path(__file__).resolve().parents[3]).as_posix()
            if workspace_root is None
            else target.name
        ),
        bytes=len(final_payload),
        sha256=final_digest,
        item_count=item_count,
    )


def prepare_document_artifact_delivery(
    result: DocumentArtifactResult,
    *,
    workspace_root: Path | None = None,
) -> DocumentArtifactDelivery:
    """Revalidate a generated artifact immediately before an external send."""
    root = _prepare_workspace(
        workspace_root or default_document_artifact_root(),
        enforce_fixed_root=workspace_root is None,
    )
    relative_file = Path(result.relative_file)
    if relative_file.is_absolute():
        raise DocumentArtifactError("delivery_path_invalid")
    if workspace_root is None:
        project_root = Path(__file__).resolve().parents[3]
        file_path = (project_root / relative_file).resolve(strict=True)
    else:
        if len(relative_file.parts) != 1:
            raise DocumentArtifactError("delivery_path_invalid")
        file_path = (root / relative_file).resolve(strict=True)
    if file_path.parent != root or file_path.is_symlink():
        raise DocumentArtifactError("delivery_path_invalid")
    extension = _COMMAND_EXTENSIONS.get(result.command)
    if extension is None or file_path.suffix.lower() != extension:
        raise DocumentArtifactError("delivery_format_invalid")
    delivery = DocumentArtifactDelivery(
        artifact_id=result.artifact_id,
        command=result.command,
        file_path=file_path,
        bytes=result.bytes,
        sha256=result.sha256,
    )
    return validate_document_artifact_delivery(
        delivery,
        workspace_root=workspace_root,
    )


def validate_document_artifact_delivery(
    delivery: DocumentArtifactDelivery,
    *,
    workspace_root: Path | None = None,
) -> DocumentArtifactDelivery:
    root = _prepare_workspace(
        workspace_root or default_document_artifact_root(),
        enforce_fixed_root=workspace_root is None,
    )
    try:
        file_path = delivery.file_path.resolve(strict=True)
    except OSError as exc:
        raise DocumentArtifactError("delivery_file_unavailable") from exc
    if file_path.parent != root or file_path.is_symlink():
        raise DocumentArtifactError("delivery_path_invalid")
    extension = _COMMAND_EXTENSIONS.get(delivery.command)
    if extension is None or file_path.suffix.lower() != extension:
        raise DocumentArtifactError("delivery_format_invalid")
    payload = file_path.read_bytes()
    digest = hashlib.sha256(payload).hexdigest()
    if len(payload) != delivery.bytes or digest != delivery.sha256:
        raise DocumentArtifactError("delivery_integrity_failed")
    if len(payload) > DOCUMENT_ARTIFACT_MAX_BYTES:
        raise DocumentArtifactError("delivery_size_invalid")
    return DocumentArtifactDelivery(
        artifact_id=delivery.artifact_id,
        command=delivery.command,
        file_path=file_path,
        bytes=len(payload),
        sha256=digest,
    )
